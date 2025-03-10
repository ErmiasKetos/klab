import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
from datetime import datetime
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Set full-window layout with expanded sidebar
st.set_page_config(layout="wide", initial_sidebar_state="expanded")

# Initialize session state for scenario management
if 'scenarios' not in st.session_state:
    st.session_state.scenarios = {}

# -------------------------------------------------------------------
# Helper Function: Calculate Monthly Capacity per Phase
# -------------------------------------------------------------------
def calculate_capacity(row, assumptions):
    base = assumptions['base_capacity']
    # Effective shift multiplier: additional shifts have half the staffing per extra shift
    effective_shift = 1 + (assumptions['shift_multiplier'] - 1) / 2
    if row['Phase'] == 'Phase 1':
        return base * assumptions['phase1_staff']
    elif row['Phase'] == 'Phase 2':
        return base * assumptions['phase2_staff'] * effective_shift
    else:  # Phase 3 with AI boost
        return base * assumptions['phase3_staff'] * effective_shift * (1 + assumptions['ai_efficiency'])

# -------------------------------------------------------------------
# Input Assumptions: Test Mix, Operational Inputs, Production, Goals, and Phase Dates
# -------------------------------------------------------------------
def input_assumptions():
    with st.sidebar:
        st.header("Model Assumptions")
        # Analyst Work Parameters
        weekly_hours = 40  # hours per week
        weeks_per_month = 4  # weeks per month
        productivity = st.slider("Analyst Productivity Factor", 0.5, 1.0, 0.8)
        
        # Test Mix Assumptions
        st.subheader("Test Mix Assumptions")
        basic_pct = st.number_input("Percentage of Basic Metals tests (%)", 0, 100, value=50)
        anions_pct = st.number_input("Percentage of Anions tests (%)", 0, 100, value=30)
        advanced_pct = st.number_input("Percentage of Advanced Metals tests (%)", 0, 100, value=20)
        total_pct = basic_pct + anions_pct + advanced_pct
        if total_pct != 100:
            st.warning(f"Test mix percentages sum to {total_pct}%. They should total 100%.")
        weighted_test_time = (basic_pct/100.0)*1 + (anions_pct/100.0)*0.5 + (advanced_pct/100.0)*1.5
        available_hours = weekly_hours * weeks_per_month * productivity
        computed_base_capacity = available_hours / weighted_test_time
        st.markdown(f"**Computed Base Capacity/Analyst/Month:** {computed_base_capacity:.0f} samples")
        
        # Core Operational Inputs
        st.subheader("Operational Inputs")
        assumptions = {
            'base_capacity': computed_base_capacity,
            'phase1_staff': st.number_input("Phase 1 Staff", 1, value=3),
            'phase2_staff': st.number_input("Phase 2 Staff", 1, value=5),
            'phase3_staff': st.number_input("Phase 3 Staff", 1, value=5),
            'shift_multiplier': st.slider("Shift Multiplier (Phase 2+)", 1.0, 3.0, 2.0),
            'ai_efficiency': st.slider("AI Efficiency Boost (Phase 3)", 0.0, 1.0, 0.3),
            'salary': st.number_input("Monthly Salary/Analyst ($)", 1000, value=5000),
            'equipment_lease': st.number_input("Equipment Lease ($/month)", 0, value=10378),
            'instrument_running': st.number_input("Instrument Running ($/month)", 0, value=2775),
            'software_licenses': st.number_input("Software/Licenses ($/month)", 0, value=2000),
            'qaqc_rate': st.slider("QA/QC Rate (% of overhead)", 0.0, 0.2, 0.08),
            'avg_test_price': st.number_input("Average Test Price ($/sample)", 1, value=300),
        }
        
        # Production Settings: Control Actual Processed Samples
        st.subheader("Production Settings")
        user_samples = st.slider("Desired Samples Processed per Month", 100, 10000, value=1000)
        assumptions['user_samples'] = user_samples
        
        # Goal Settings: Option for Monthly Sample or Profit Goal
        st.subheader("Goal Settings")
        goal_type = st.radio("Select Goal Type", ["Monthly Sample Goal", "Monthly Profit Goal"])
        if goal_type == "Monthly Sample Goal":
            goal_value = st.number_input("Monthly Sample Goal (tests/month)", 100, 10000, value=1000)
        else:
            goal_value = st.number_input("Monthly Profit Goal ($/month)", 1000, value=50000)
        assumptions['goal_type'] = goal_type
        assumptions['goal_value'] = goal_value
        
        # Phase Dates
        st.subheader("Phase Dates")
        assumptions['phase1_start'] = pd.to_datetime(st.date_input("Phase 1 Start", datetime(2025, 4, 1)))
        assumptions['phase2_start'] = pd.to_datetime(st.date_input("Phase 2 Start", datetime(2026, 1, 1)))
        assumptions['phase3_start'] = pd.to_datetime(st.date_input("Phase 3 Start", datetime(2027, 1, 1)))
        assumptions['phase3_end']   = pd.to_datetime(st.date_input("Model End Date", datetime(2027, 12, 31)))
        
        return assumptions

# -------------------------------------------------------------------
# Timeline Generation: Calculate Monthly Metrics and Goal Achievement
# -------------------------------------------------------------------
def generate_timeline(assumptions):
    dates = pd.date_range(start=assumptions['phase1_start'],
                          end=assumptions['phase3_end'],
                          freq='MS')
    timeline = pd.DataFrame(index=dates)
    timeline.index.name = 'Month'
    
    # Assign Phases
    timeline['Phase'] = 'Phase 1'
    timeline.loc[timeline.index >= assumptions['phase2_start'], 'Phase'] = 'Phase 2'
    timeline.loc[timeline.index >= assumptions['phase3_start'], 'Phase'] = 'Phase 3'
    
    # Theoretical Monthly Capacity
    timeline['Monthly Capacity'] = timeline.apply(lambda x: calculate_capacity(x, assumptions), axis=1)
    timeline['Cumulative Samples'] = timeline['Monthly Capacity'].cumsum()
    
    # Processed Samples: Minimum of computed capacity and user-controlled production
    timeline['Processed Samples'] = timeline['Monthly Capacity'].apply(
        lambda cap: min(cap, assumptions.get('user_samples', cap))
    )
    
    # Cost Calculations
    timeline['Staff Costs'] = timeline['Phase'].map({
        'Phase 1': assumptions['phase1_staff'] * assumptions['salary'],
        'Phase 2': assumptions['phase2_staff'] * assumptions['salary'],
        'Phase 3': assumptions['phase3_staff'] * assumptions['salary']
    })
    equipment_lease = assumptions.get('equipment_lease', 10378)
    instrument_running = assumptions.get('instrument_running', 2775)
    software_licenses = assumptions.get('software_licenses', 2000)
    qaqc_rate = assumptions.get('qaqc_rate', 0.08)
    overhead_base = equipment_lease + instrument_running + software_licenses
    timeline['Overhead Costs'] = overhead_base + overhead_base * qaqc_rate
    timeline['Total Cost'] = timeline['Staff Costs'] + timeline['Overhead Costs']
    
    # Revenue is based on Processed Samples
    avg_test_price = assumptions.get('avg_test_price', 300)
    timeline['Revenue'] = timeline['Processed Samples'] * avg_test_price
    timeline['Profit'] = timeline['Revenue'] - timeline['Total Cost']
    
    # Determine "Goal Met?" based on selected goal type
    if assumptions['goal_type'] == "Monthly Sample Goal":
        timeline['Goal Met?'] = timeline['Processed Samples'].apply(
            lambda x: "Goal Met" if x >= assumptions['goal_value'] else "Under Target"
        )
    else:
        timeline['Goal Met?'] = timeline['Profit'].apply(
            lambda x: "Goal Met" if x >= assumptions['goal_value'] else "Under Target"
        )
    
    return timeline

# -------------------------------------------------------------------
# Render Base Visualizations: Graphs and Detailed Data Table
# -------------------------------------------------------------------
def render_base_visualizations(timeline):
    st.subheader("Monthly Capacity")
    fig_cap = px.line(timeline.reset_index(), x='Month', y='Monthly Capacity', title="Theoretical Monthly Capacity Over Time")
    st.plotly_chart(fig_cap, use_container_width=True)
    
    st.subheader("Processed Samples")
    fig_proc = px.line(timeline.reset_index(), x='Month', y='Processed Samples', title="Actual Samples Processed Over Time")
    st.plotly_chart(fig_proc, use_container_width=True)
    
    st.subheader("Monthly Revenue")
    fig_rev = px.line(timeline.reset_index(), x='Month', y='Revenue', title="Monthly Revenue Over Time")
    st.plotly_chart(fig_rev, use_container_width=True)
    
    st.subheader("Monthly Profit")
    fig_profit = px.line(timeline.reset_index(), x='Month', y='Profit', title="Monthly Profit Over Time")
    st.plotly_chart(fig_profit, use_container_width=True)
    
    st.subheader("Detailed Data")
    st.dataframe(timeline[["Phase", "Monthly Capacity", "Processed Samples", "Goal Met?", "Revenue", "Total Cost", "Profit"]])

# -------------------------------------------------------------------
# Scenario Management: Save, Load, and Delete Scenarios
# -------------------------------------------------------------------
def scenario_management(assumptions):
    with st.sidebar:
        st.subheader("Scenario Management")
        scenario_name = st.text_input("Scenario Name")
        col1, col2 = st.columns(2)
        with col1:
            if st.button("💾 Save Current Settings"):
                if scenario_name:
                    st.session_state.scenarios[scenario_name] = assumptions.copy()
                    st.success(f"Scenario '{scenario_name}' saved!")
        selected = st.multiselect("Compare Scenarios", options=list(st.session_state.scenarios.keys()),
                                    default=[list(st.session_state.scenarios.keys())[-1]] if st.session_state.scenarios else [])
        with col2:
            if st.button("🗑️ Delete Selected"):
                for name in selected:
                    if name in st.session_state.scenarios:
                        del st.session_state.scenarios[name]
        return selected

# -------------------------------------------------------------------
# Generate PPT for Current Scenario (Beautiful and Professional)
# -------------------------------------------------------------------
def generate_ppt(assumptions, timeline):
    prs = Presentation()
    
    # Helper: Add a border to a slide
    def add_border(slide):
        left = Inches(0.2)
        top = Inches(0.2)
        width = prs.slide_width - Inches(0.4)
        height = prs.slide_height - Inches(0.4)
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        border.fill.solid()
        border.fill.fore_color.rgb = RGBColor(255, 255, 255)  # white fill
        border.line.color.rgb = RGBColor(0, 51, 102)  # deep blue
        border.line.width = Pt(3)
        # Send border to back
        border._element.getparent().insert(0, border._element)
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    add_border(slide)
    slide.shapes.title.text = "Current Scenario: Lab Scalability & Profitability"
    slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    # Assumptions Slide with Card Effect
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_border(slide)
    # Create a card (rounded rectangle)
    card_left = Inches(0.5)
    card_top = Inches(0.5)
    card_width = prs.slide_width - Inches(1)
    card_height = prs.slide_height - Inches(1)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(240, 240, 240)  # light grey
    card.line.color.rgb = RGBColor(0, 51, 102)
    card.line.width = Pt(2)
    # Add assumptions text on top of card
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(8.6), Inches(6))
    tf = txBox.text_frame
    tf.text = "Model Assumptions:\n"
    for key, value in assumptions.items():
        tf.text += f"{key}: {value}\n"
    
    # Create a slide for each graph
    figures = [
        ("Theoretical Monthly Capacity", px.line(timeline.reset_index(), x='Month', y='Monthly Capacity', title="Theoretical Monthly Capacity")),
        ("Processed Samples", px.line(timeline.reset_index(), x='Month', y='Processed Samples', title="Processed Samples")),
        ("Monthly Revenue", px.line(timeline.reset_index(), x='Month', y='Revenue', title="Monthly Revenue")),
        ("Monthly Profit", px.line(timeline.reset_index(), x='Month', y='Profit', title="Monthly Profit"))
    ]
    
    for title_text, fig in figures:
        slide = prs.slides.add_slide(blank_slide_layout)
        add_border(slide)
        # Add slide title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        tf.text = title_text
        # Export Plotly figure as PNG (requires kaleido)
        img_bytes = fig.to_image(format="png")
        image_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(image_stream, Inches(1), Inches(1.2), width=Inches(8))
    
    # Detailed Data Slide (first 10 rows)
    slide = prs.slides.add_slide(blank_slide_layout)
    add_border(slide)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
    tf = txBox.text_frame
    tf.text = "Detailed Data (first 10 rows):\n"
    df_str = timeline.head(10).to_string()
    tf.text += df_str
    
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# -------------------------------------------------------------------
# Render Scenario Comparison & KPI Dashboard
# -------------------------------------------------------------------
def render_comparison(selected_scenarios):
    valid_scenarios = [name for name in selected_scenarios if name in st.session_state.scenarios]
    if len(valid_scenarios) < 1:
        return
    
    all_data = []
    metrics = []
    for name in valid_scenarios:
        scenario_assumptions = st.session_state.scenarios[name]
        timeline = generate_timeline(scenario_assumptions)
        timeline['Scenario'] = name
        all_data.append(timeline)
        
        start_date = timeline.index[0]
        if scenario_assumptions['goal_type'] == "Monthly Sample Goal":
            goal_achieved = timeline[timeline['Processed Samples'] >= scenario_assumptions['goal_value']]
        else:
            goal_achieved = timeline[timeline['Profit'] >= scenario_assumptions['goal_value']]
        
        if not goal_achieved.empty:
            goal_date = goal_achieved.index[0]
            months_to_goal = (goal_date.year - start_date.year) * 12 + (goal_date.month - start_date.month)
            goal_date_str = goal_date.strftime("%Y-%m-%d")
        else:
            goal_date_str = "Not Reached"
            months_to_goal = "Not Reached"
        
        total_months = len(timeline)
        total_revenue = timeline['Revenue'].sum()
        total_cost = timeline['Total Cost'].sum()
        total_profit = timeline['Profit'].sum()
        profit_margin = (total_profit / total_revenue * 100) if total_revenue else 0
        avg_monthly_revenue = total_revenue / total_months
        avg_monthly_profit = total_profit / total_months
        total_samples = timeline['Cumulative Samples'].iloc[-1]
        cost_per_test = total_cost / total_samples if total_samples > 0 else None
        
        metrics.append({
            'Scenario': name,
            'Peak Capacity (tests/month)': timeline['Monthly Capacity'].max(),
            'Peak Processed (tests/month)': timeline['Processed Samples'].max(),
            'Goal Achieved Date': goal_date_str,
            'Months to Reach Goal': months_to_goal,
            'Total Revenue ($M)': total_revenue / 1e6,
            'Total Cost ($M)': total_cost / 1e6,
            'Total Profit ($M)': total_profit / 1e6,
            'Profit Margin (%)': profit_margin,
            'Avg Monthly Revenue ($K)': avg_monthly_revenue / 1e3,
            'Avg Monthly Profit ($K)': avg_monthly_profit / 1e3,
            'Cost per Test ($)': cost_per_test
        })
    
    df_combined = pd.concat(all_data)
    
    st.subheader("Scenario Comparison Chart")
    first_scenario = st.session_state.scenarios[valid_scenarios[0]]
    goal_value = first_scenario.get('goal_value', 1000)
    if first_scenario['goal_type'] == "Monthly Sample Goal":
        fig = px.line(df_combined.reset_index(), x='Month', y='Processed Samples', color='Scenario',
                      title="Processed Samples Across Scenarios")
        fig.add_hline(y=goal_value, line_dash="dot", line_color="red")
    else:
        fig = px.line(df_combined.reset_index(), x='Month', y='Profit', color='Scenario',
                      title="Monthly Profit Across Scenarios")
        fig.add_hline(y=goal_value, line_dash="dot", line_color="red")
    st.plotly_chart(fig, use_container_width=True)
    
    st.subheader("Scenario Comparison - Profit")
    fig_profit = px.line(df_combined.reset_index(), x='Month', y='Profit', color='Scenario',
                         title="Monthly Profit Across Scenarios")
    st.plotly_chart(fig_profit, use_container_width=True)
    
    st.subheader("Key Metrics (KPIs)")
    df_metrics = pd.DataFrame(metrics)
    st.dataframe(df_metrics, use_container_width=True)
    
    st.markdown("### KPI Explanations:")
    st.markdown("""
    - **Peak Capacity (tests/month):** The maximum theoretical tests that could be processed in a month.
    - **Peak Processed (tests/month):** The maximum number of samples actually processed (user-controlled) in a month.
    - **Goal Achieved Date:** The first month when processed samples or profit meets/exceeds the set goal.
    - **Months to Reach Goal:** Number of months from start until the goal is achieved.
    - **Total Revenue ($M):** Cumulative revenue over the period (in millions), calculated using actual processed samples.
    - **Total Cost ($M):** Cumulative cost over the period (in millions).
    - **Total Profit ($M):** Total revenue minus total cost (in millions).
    - **Profit Margin (%):** Total profit as a percentage of total revenue.
    - **Avg Monthly Revenue ($K):** Average monthly revenue (in thousands).
    - **Avg Monthly Profit ($K):** Average monthly profit (in thousands).
    - **Cost per Test ($):** Average cost incurred per test performed.
    """)
    
# -------------------------------------------------------------------
# Monte Carlo Simulation & Interpretation
# -------------------------------------------------------------------
def run_monte_carlo_simulation(base_assumptions, n_simulations=500):
    results = []
    progress_bar = st.progress(0)
    for i in range(n_simulations):
        perturbed = base_assumptions.copy()
        perturbed['base_capacity'] = np.random.normal(
            base_assumptions['base_capacity'],
            base_assumptions['base_capacity'] * 0.1
        )
        perturbed['ai_efficiency'] = np.random.triangular(
            left=0.1, mode=base_assumptions['ai_efficiency'], right=0.5
        )
        perturbed['phase2_staff'] = int(np.random.choice([
            base_assumptions['phase2_staff'] - 1,
            base_assumptions['phase2_staff'],
            base_assumptions['phase2_staff'] + 1
        ]))
        timeline = generate_timeline(perturbed)
        if base_assumptions['goal_type'] == "Monthly Sample Goal":
            goal_met = timeline[timeline['Processed Samples'] >= base_assumptions['goal_value']]
        else:
            goal_met = timeline[timeline['Profit'] >= base_assumptions['goal_value']]
        results.append({
            'peak_capacity': timeline['Monthly Capacity'].max(),
            'months_to_goal': ((goal_met.index[0] - timeline.index[0]).days // 30) if not goal_met.empty else None,
            'total_cost': timeline['Total Cost'].sum() / 1e6
        })
        progress_bar.progress((i + 1) / n_simulations)
    return pd.DataFrame(results)

def render_monte_carlo(base_assumptions):
    st.subheader("Risk Analysis (Monte Carlo Simulation)")
    with st.expander("⚙️ Simulation Settings"):
        n_simulations = st.number_input("Number of Simulations", 100, 5000, 500)
        run_sim = st.button("Run Simulation")
    if run_sim:
        with st.spinner(f"Running {n_simulations} simulations..."):
            results = run_monte_carlo_simulation(base_assumptions, n_simulations)
        col1, col2, col3 = st.columns(3)
        with col1:
            fig1 = px.histogram(results, x='peak_capacity', title="Peak Capacity Distribution")
            st.plotly_chart(fig1, use_container_width=True)
        with col2:
            valid_times = results[results['months_to_goal'].notna()]
            fig2 = px.histogram(valid_times, x='months_to_goal', title="Time to Reach Goal (Months)")
            st.plotly_chart(fig2, use_container_width=True)
        with col3:
            fig3 = px.scatter(results, x='total_cost', y='peak_capacity', title="Cost vs Capacity Tradeoff")
            st.plotly_chart(fig3, use_container_width=True)
        avg_peak = results['peak_capacity'].mean()
        median_peak = results['peak_capacity'].median()
        std_peak = results['peak_capacity'].std()
        valid_times = results[results['months_to_goal'].notna()]
        if not valid_times.empty:
            avg_time = valid_times['months_to_goal'].mean()
            median_time = valid_times['months_to_goal'].median()
        else:
            avg_time = None
            median_time = None
        probability_goal = (1 - results['months_to_goal'].isna().mean()) * 100
        st.metric("Probability of Achieving Goal", f"{probability_goal:.1f}%")
        st.markdown("### Monte Carlo Simulation Interpretation")
        st.markdown(
            f"**Peak Capacity Distribution:** The simulations show an average peak capacity of **{avg_peak:.1f} tests/month** "
            f"(median: **{median_peak:.1f}**, std: **{std_peak:.1f}**), reflecting variability in operational performance."
        )
        if avg_time is not None:
            st.markdown(
                f"**Time to Reach Goal:** For runs that met the goal, the average time was **{avg_time:.1f} months** "
                f"(median: **{median_time:.1f} months**)."
            )
        else:
            st.markdown("**Time to Reach Goal:** The goal was not reached in most simulations.")
        st.markdown(
            "**Cost vs Capacity Tradeoff:** Generally, higher capacity is associated with higher costs, highlighting a key tradeoff when scaling operations."
        )
        st.markdown(
            f"**Probability of Achieving Goal:** There is a **{probability_goal:.1f}%** chance of meeting the goal "
            f"({base_assumptions['goal_type']} of {base_assumptions['goal_value']})."
        )

# -------------------------------------------------------------------
# Main Application Structure
# -------------------------------------------------------------------
def main():
    st.title("KELAB Scalability & Profitability Modeling")
    st.markdown("""
    This app models lab performance from **April 2025** to **December 2027**.
    Set your operational goals (either as a Monthly Sample Goal or Monthly Profit Goal) and adjust production via the slider.
    You can also generate a downloadable PPT file for the Current Scenario (including graphs, assumptions, and detailed data).
    """)
    
    assumptions = input_assumptions()
    selected_scenarios = scenario_management(assumptions)
    
    tab1, tab2, tab3 = st.tabs(["Current Scenario", "Scenario Comparison", "Risk Analysis"])
    
    with tab1:
        timeline = generate_timeline(assumptions)
        render_base_visualizations(timeline)
        
        st.markdown("## Download Current Scenario as PPT")
        if st.button("Generate PPT"):
            ppt_file = generate_ppt(assumptions, timeline)
            st.download_button(label="Download PPT",
                               data=ppt_file,
                               file_name="Current_Scenario.pptx",
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    
    with tab2:
        render_comparison(selected_scenarios)
    
    with tab3:
        render_monte_carlo(assumptions)

if __name__ == "__main__":
    main()
